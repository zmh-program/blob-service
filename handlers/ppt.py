from io import BytesIO

from fastapi import UploadFile
import pptx


def is_pptx(filename: str) -> bool:
    """Return True if file is a PowerPoint presentation."""
    return filename.endswith(".pptx") or filename.endswith(".ppt")


def process(file: UploadFile) -> str:
    """Process PowerPoint presentation and return its contents."""
    if file.filename.endswith(".ppt"):
        raise ValueError(".ppt files are not supported, only .pptx files are supported.")

    content = file.file.read()
    return pptx_to_text(content)


def pptx_to_text(content: bytes) -> str:
    """Convert PowerPoint presentation to text."""
    # fix: bytes object has no attribute 'seek'

    prs = pptx.Presentation(BytesIO(content))
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)

